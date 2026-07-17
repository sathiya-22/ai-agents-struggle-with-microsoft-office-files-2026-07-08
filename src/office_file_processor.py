import docx
from openpyxl import load_workbook
from pptx import Presentation
import os
import requests

class OfficeFileProcessor:
    def __init__(self, api_key):
        self.api_key = api_key

    def extract_text_from_word(self, file_path):
        doc = docx.Document(file_path)
        text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        return text

    def extract_text_from_excel(self, file_path):
        wb = load_workbook(file_path)
        text = ''
        for sheet in wb.worksheets:
            for row in sheet.rows:
                for cell in row:
                    if cell.value is not None:
                        text += str(cell.value) + ' '
        return text

    def extract_text_from_pptx(self, file_path):
        presentation = Presentation(file_path)
        text = ''
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.has_text:
                    text += shape.text_frame.text + ' '
        return text

    def process_office_file(self, file_path):
        file_extension = os.path.splitext(file_path)[1].lower()
        extracted_text = ""
        if file_extension == '.docx':
            extracted_text = self.extract_text_from_word(file_path)
        elif file_extension == '.xlsx':
            extracted_text = self.extract_text_from_excel(file_path)
        elif file_extension == '.pptx':
            extracted_text = self.extract_text_from_pptx(file_path)
        else:
            raise ValueError("Unsupported file type. Only .docx, .xlsx, and .pptx are supported.")
        return extracted_text

    def send_to_geminai(self, text):
        # NOTE: This URL is a placeholder. The actual Gemini API URL and request body
        # would need to be updated based on the specific Gemini API documentation.
        url = 'https://api.gemini.ai/v1/analyze' # This is a hypothetical endpoint
        headers = {'Authorization': f'Bearer {self.api_key}', 'Content-Type': 'application/json'}
        # A more realistic request for a text analysis API might look like this:
        payload = {
            "contents": [
                {
                    "parts": [
                        {"text": text}
                    ]
                }
            ]
        }
        response = requests.post(url, headers=headers, json=payload)
        response.raise_for_status() # Raise an exception for bad status codes
        return response.json()
